import sys
import os
import subprocess

# 1. Dependency check and auto-installation
try:
    import pptx
except ImportError:
    print("WARNING: python-pptx is required. Attempting to install it now...")
    try:
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
        import pptx
        print("python-pptx installed successfully!")
    except Exception as e:
        print(f"ERROR: Failed to install python-pptx automatically: {e}")
        print("Please run: pip install python-pptx")
        sys.exit(1)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # Set slide dimensions to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Add a blank slide
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # 2. Set slide background to Deep Dark Indigo #0b132b
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(11, 19, 43)
    
    # 3. Add Slide Header
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf_header = header_box.text_frame
    tf_header.word_wrap = True
    tf_header.margin_left = tf_header.margin_right = tf_header.margin_top = tf_header.margin_bottom = 0
    
    p_header = tf_header.paragraphs[0]
    # Brand: AQS Inspect
    run_brand = p_header.add_run()
    run_brand.text = "AQS Inspect "
    run_brand.font.name = "Arial"
    run_brand.font.size = Pt(22)
    run_brand.font.bold = True
    run_brand.font.color.rgb = RGBColor(6, 182, 212) # Electric Blue
    
    # Title Details
    run_title = p_header.add_run()
    run_title.text = "— AI-Driven Review & Compliance Transformation"
    run_title.font.name = "Arial"
    run_title.font.size = Pt(20)
    run_title.font.bold = True
    run_title.font.color.rgb = RGBColor(241, 245, 249) # Clean White
    
    # Subtitle
    p_sub = tf_header.add_paragraph()
    p_sub.text = "Standardizing enterprise source code, enforcing AQS Guidelines, and accelerating pull request delivery."
    p_sub.font.name = "Arial"
    p_sub.font.size = Pt(10.5)
    p_sub.font.color.rgb = RGBColor(148, 163, 184) # Muted gray
    p_sub.space_before = Pt(4)
    
    # 4. Top Grid: 3 Columns
    # Widths and positions
    card_width = Inches(3.9)
    card_height = Inches(3.7)
    card_y = Inches(1.3)
    
    col_x_coords = [Inches(0.5), Inches(4.716), Inches(8.933)]
    
    card_bg_color = RGBColor(30, 41, 59)     # Dark card background #1e293b
    card_border_color = RGBColor(51, 65, 85) # Card border #334155
    
    # Accent colors for the left bars
    red_accent = RGBColor(244, 63, 94)      # Blocker Red
    blue_accent = RGBColor(6, 182, 212)     # Electric Blue
    green_accent = RGBColor(16, 185, 129)    # Emerald Green
    
    card_accents = [red_accent, blue_accent, green_accent]
    
    # Draw backgrounds, left accent bars, and outline borders
    for i, x in enumerate(col_x_coords):
        # Base Card Rounded Rectangle
        card_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, card_y, card_width, card_height
        )
        card_shape.fill.solid()
        card_shape.fill.fore_color.rgb = card_bg_color
        card_shape.line.color.rgb = card_border_color
        card_shape.line.width = Pt(1)
        
        # Left Accent Border Strip
        accent_strip = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x + Inches(0.01), card_y + Inches(0.04), Inches(0.06), card_height - Inches(0.08)
        )
        accent_strip.fill.solid()
        accent_strip.fill.fore_color.rgb = card_accents[i]
        accent_strip.line.fill.background() # borderless
        
    # --- Column 1 text: The Bottleneck ---
    tx_col1 = slide.shapes.add_textbox(col_x_coords[0] + Inches(0.2), card_y + Inches(0.15), card_width - Inches(0.3), card_height - Inches(0.3))
    tf1 = tx_col1.text_frame
    tf1.word_wrap = True
    
    # Title
    p_t1 = tf1.paragraphs[0]
    p_t1.text = "THE CODE REVIEW BOTTLENECK"
    p_t1.font.name = "Arial"
    p_t1.font.size = Pt(12.5)
    p_t1.font.bold = True
    p_t1.font.color.rgb = red_accent
    p_t1.space_after = Pt(12)
    
    # Bottleneck list items
    bottlenecks = [
        ("4+ Hours", "Manual Effort", "Developers and leads spend substantial manual effort on code checks per PR."),
        ("Expert Heavy", "Lead Dependency", "Senior resources spend valuable time validating style rules and guidelines."),
        ("Inconsistent", "Variable Quality", "Manual code audits regularly miss compliance guidelines and logic defects."),
        ("Slow Cycle", "Delayed Delivery", "Unreviewed pull request backlogs stack up, delaying time-to-market.")
    ]
    
    for label, subtitle, desc in bottlenecks:
        p_item = tf1.add_paragraph()
        p_item.space_after = Pt(8)
        
        # Prefix label
        run_lbl = p_item.add_run()
        run_lbl.text = f"[{label}] "
        run_lbl.font.name = "Arial"
        run_lbl.font.size = Pt(9.5)
        run_lbl.font.bold = True
        run_lbl.font.color.rgb = red_accent
        
        # Subtitle
        run_sub = p_item.add_run()
        run_sub.text = f"{subtitle}: "
        run_sub.font.name = "Arial"
        run_sub.font.size = Pt(9.5)
        run_sub.font.bold = True
        run_sub.font.color.rgb = RGBColor(241, 245, 249)
        
        # Description
        run_desc = p_item.add_run()
        run_desc.text = desc
        run_desc.font.name = "Arial"
        run_desc.font.size = Pt(9)
        run_desc.font.color.rgb = RGBColor(148, 163, 184)

    # --- Column 2 text: Time Transformation ---
    tx_col2 = slide.shapes.add_textbox(col_x_coords[1] + Inches(0.2), card_y + Inches(0.15), card_width - Inches(0.3), card_height - Inches(0.3))
    tf2 = tx_col2.text_frame
    tf2.word_wrap = True
    
    # Title
    p_t2 = tf2.paragraphs[0]
    p_t2.text = "DRAMATIC TIME TRANSFORMATION"
    p_t2.font.name = "Arial"
    p_t2.font.size = Pt(12.5)
    p_t2.font.bold = True
    p_t2.font.color.rgb = blue_accent
    p_t2.space_after = Pt(18)
    
    # Manual vs AQS Inspect block layout
    # We will position separate text boxes for Manual (Left), Arrow (Center), and AQS (Right)
    inner_y = card_y + Inches(0.8)
    
    # Manual box
    manual_box = slide.shapes.add_textbox(col_x_coords[1] + Inches(0.2), inner_y, Inches(1.3), Inches(1.2))
    tf_m = manual_box.text_frame
    tf_m.word_wrap = True
    p_ml = tf_m.paragraphs[0]
    p_ml.text = "MANUAL"
    p_ml.font.name = "Arial"
    p_ml.font.size = Pt(9)
    p_ml.font.bold = True
    p_ml.font.color.rgb = RGBColor(148, 163, 184)
    p_ml.alignment = PP_ALIGN.CENTER
    
    p_mv = tf_m.add_paragraph()
    p_mv.text = "4.0"
    p_mv.font.name = "Arial"
    p_mv.font.size = Pt(38)
    p_mv.font.bold = True
    p_mv.font.color.rgb = red_accent
    p_mv.alignment = PP_ALIGN.CENTER
    
    p_mu = tf_m.add_paragraph()
    p_mu.text = "hours"
    p_mu.font.name = "Arial"
    p_mu.font.size = Pt(9.5)
    p_mu.font.color.rgb = RGBColor(148, 163, 184)
    p_mu.alignment = PP_ALIGN.CENTER
    
    # Arrow box
    arrow_box = slide.shapes.add_textbox(col_x_coords[1] + Inches(1.5), inner_y + Inches(0.2), Inches(0.9), Inches(0.8))
    tf_a = arrow_box.text_frame
    p_a = tf_a.paragraphs[0]
    p_a.text = "➔"
    p_a.font.name = "Arial"
    p_a.font.size = Pt(32)
    p_a.font.color.rgb = blue_accent
    p_a.alignment = PP_ALIGN.CENTER
    
    # AQS Box
    aqs_box = slide.shapes.add_textbox(col_x_coords[1] + Inches(2.4), inner_y, Inches(1.3), Inches(1.2))
    tf_aq = aqs_box.text_frame
    tf_aq.word_wrap = True
    p_aql = tf_aq.paragraphs[0]
    p_aql.text = "AQS INSPECT"
    p_aql.font.name = "Arial"
    p_aql.font.size = Pt(9)
    p_aql.font.bold = True
    p_aql.font.color.rgb = RGBColor(148, 163, 184)
    p_aql.alignment = PP_ALIGN.CENTER
    
    p_aqv = tf_aq.add_paragraph()
    p_aqv.text = "~5"
    p_aqv.font.name = "Arial"
    p_aqv.font.size = Pt(38)
    p_aqv.font.bold = True
    p_aqv.font.color.rgb = green_accent
    p_aqv.alignment = PP_ALIGN.CENTER
    
    p_aqu = tf_aq.add_paragraph()
    p_aqu.text = "minutes"
    p_aqu.font.name = "Arial"
    p_aqu.font.size = Pt(9.5)
    p_aqu.font.color.rgb = RGBColor(148, 163, 184)
    p_aqu.alignment = PP_ALIGN.CENTER
    
    # Inner Banner Blue Box
    banner_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        col_x_coords[1] + Inches(0.2), card_y + Inches(2.3), card_width - Inches(0.4), Inches(1.1)
    )
    banner_shape.fill.solid()
    banner_shape.fill.fore_color.rgb = RGBColor(29, 78, 216) # blue-700
    banner_shape.line.fill.background()
    
    tx_banner = slide.shapes.add_textbox(col_x_coords[1] + Inches(0.2), card_y + Inches(2.3), card_width - Inches(0.4), Inches(1.1))
    tf_b = tx_banner.text_frame
    tf_b.word_wrap = True
    
    p_bm = tf_b.paragraphs[0]
    p_bm.text = "98% TIME REDUCTION"
    p_bm.font.name = "Arial"
    p_bm.font.size = Pt(16)
    p_bm.font.bold = True
    p_bm.font.color.rgb = RGBColor(255, 255, 255)
    p_bm.alignment = PP_ALIGN.CENTER
    p_bm.space_after = Pt(2)
    p_bm.space_before = Pt(8)
    
    p_bd = tf_b.add_paragraph()
    p_bd.text = "Automated compliance audits & LLM review automation"
    p_bd.font.name = "Arial"
    p_bd.font.size = Pt(9)
    p_bd.font.color.rgb = RGBColor(147, 197, 253) # blue-300
    p_bd.alignment = PP_ALIGN.CENTER

    # --- Column 3 text: Output Quality Comparison Table ---
    tx_col3 = slide.shapes.add_textbox(col_x_coords[2] + Inches(0.2), card_y + Inches(0.15), card_width - Inches(0.3), Inches(0.6))
    tf3 = tx_col3.text_frame
    p_t3 = tf3.paragraphs[0]
    p_t3.text = "QUALITY & COMPLIANCE TRANSFORMATION"
    p_t3.font.name = "Arial"
    p_t3.font.size = Pt(12.5)
    p_t3.font.bold = True
    p_t3.font.color.rgb = green_accent
    
    # Table details
    rows = 6
    cols = 3
    t_left = col_x_coords[2] + Inches(0.15)
    t_top = card_y + Inches(0.7)
    t_width = card_width - Inches(0.3)
    t_height = Inches(2.7)
    
    table_shape = slide.shapes.add_table(rows, cols, t_left, t_top, t_width, t_height)
    table = table_shape.table
    
    # Set column widths
    table.columns[0].width = Inches(1.2)
    table.columns[1].width = Inches(1.1)
    table.columns[2].width = Inches(1.3)
    
    # Table data
    table_data = [
        ["PARAMETER", "MANUAL", "AQS INSPECT"],
        ["Defect Detect", "Variable / Late", "90% Bugs Early"],
        ["Standardization", "Developer-dept", "100% Adherence"],
        ["Guidelines", "Variable", "Mandatory Checked"],
        ["IFS ERP Risks", "Manual Check", "Automated Hook"],
        ["Traceability", "None / Informal", "Full Audit Trail"]
    ]
    
    for r_idx in range(rows):
        for c_idx in range(cols):
            cell = table.cell(r_idx, c_idx)
            cell.text = table_data[r_idx][c_idx]
            
            # Formatting
            cell.fill.solid()
            # Set cell bg to slightly darker than card to contrast
            cell.fill.fore_color.rgb = RGBColor(24, 33, 47)
            
            p_cell = cell.text_frame.paragraphs[0]
            p_cell.font.name = "Arial"
            p_cell.alignment = PP_ALIGN.LEFT
            
            if r_idx == 0:
                p_cell.font.size = Pt(8.5)
                p_cell.font.bold = True
                p_cell.font.color.rgb = RGBColor(148, 163, 184) # light gray
            else:
                p_cell.font.size = Pt(8)
                if c_idx == 0:
                    p_cell.font.bold = True
                    p_cell.font.color.rgb = RGBColor(226, 232, 240) # clean white
                elif c_idx == 1:
                    p_cell.font.color.rgb = RGBColor(148, 163, 184) # muted
                elif c_idx == 2:
                    p_cell.font.bold = True
                    p_cell.font.color.rgb = RGBColor(56, 189, 248) # electric blue-300
                    
    # 5. Bottom Section: Development Workflow timeline
    wf_y = Inches(5.2)
    wf_width = Inches(12.333)
    wf_height = Inches(1.8)
    
    # Container box
    wf_container = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), wf_y, wf_width, wf_height
    )
    wf_container.fill.solid()
    wf_container.fill.fore_color.rgb = RGBColor(22, 28, 45) # Very dark navy #161c2d
    wf_container.line.color.rgb = card_border_color
    wf_container.line.width = Pt(1)
    
    # Section Header
    tx_wf_title = slide.shapes.add_textbox(Inches(0.7), wf_y + Inches(0.1), Inches(11.9), Inches(0.4))
    tf_wft = tx_wf_title.text_frame
    p_wft = tf_wft.paragraphs[0]
    p_wft.text = "HOW AQS INSPECT INTEGRATES INTO DEVELOPMENT WORKFLOW"
    p_wft.font.name = "Arial"
    p_wft.font.size = Pt(11)
    p_wft.font.bold = True
    p_wft.font.color.rgb = RGBColor(241, 245, 249)
    
    # Steps timeline layout
    # Width of step card = 1.7 inches, connector = 0.3 inches
    step_width = Inches(1.7)
    step_y = wf_y + Inches(0.5)
    step_h = Inches(1.1)
    
    wf_steps = [
        ("1", "PR Submission", "Developer creates a PR in Git/Azure DevOps."),
        ("2", "Context Fetch", "AQS Inspect fetches diff & ERP metadata."),
        ("3", "Standard scan", "Rule engine validates AQS Guidelines."),
        ("4", "AI Compliance", "LLM reviews edge cases and logic paths."),
        ("5", "Audit & Merge", "Inline fixes applied; audit event logged.")
    ]
    
    start_x_wf = Inches(0.9)
    gap_x_wf = Inches(2.1) # 1.7 + 0.4
    
    # Render steps and connectors
    for idx, (num, title, desc) in enumerate(wf_steps):
        current_x = start_x_wf + (idx * gap_x_wf)
        
        # Step shape
        s_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            current_x, step_y, step_width, step_h
        )
        s_shape.fill.solid()
        s_shape.fill.fore_color.rgb = RGBColor(15, 23, 42)
        # Highlight active step
        if idx == 0:
            s_shape.line.color.rgb = blue_accent
        else:
            s_shape.line.color.rgb = card_border_color
        s_shape.line.width = Pt(1)
        
        # Add text
        tx_step = slide.shapes.add_textbox(current_x + Inches(0.05), step_y + Inches(0.05), step_width - Inches(0.1), step_h - Inches(0.1))
        tf_s = tx_step.text_frame
        tf_s.word_wrap = True
        tf_s.margin_left = tf_s.margin_right = tf_s.margin_top = tf_s.margin_bottom = Inches(0.02)
        
        p_num = tf_s.paragraphs[0]
        p_num.text = f"{num}."
        p_num.font.name = "Arial"
        p_num.font.size = Pt(8.5)
        p_num.font.bold = True
        p_num.font.color.rgb = blue_accent
        
        p_st = tf_s.add_paragraph()
        p_st.text = title
        p_st.font.name = "Arial"
        p_st.font.size = Pt(9)
        p_st.font.bold = True
        p_st.font.color.rgb = RGBColor(255, 255, 255)
        p_st.space_after = Pt(2)
        
        p_sd = tf_s.add_paragraph()
        p_sd.text = desc
        p_sd.font.name = "Arial"
        p_sd.font.size = Pt(7.5)
        p_sd.font.color.rgb = RGBColor(148, 163, 184)
        
        # Draw arrow connector
        if idx < 4:
            arrow_x = current_x + step_width + Inches(0.08)
            tx_ar = slide.shapes.add_textbox(arrow_x, step_y + Inches(0.35), Inches(0.24), Inches(0.4))
            tf_ar = tx_ar.text_frame
            p_ar = tf_ar.paragraphs[0]
            p_ar.text = "➔"
            p_ar.font.name = "Arial"
            p_ar.font.size = Pt(12)
            p_ar.font.color.rgb = RGBColor(71, 85, 105) # slate-600
            p_ar.alignment = PP_ALIGN.CENTER
            
    # Save presentation
    output_path = "AQS_Inspect_Workflow_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {os.path.abspath(output_path)}")

if __name__ == "__main__":
    create_presentation()
